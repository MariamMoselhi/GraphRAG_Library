"""
PPTX Extractor for HoloLearn
Extracts text content from PowerPoint presentations.
"""

from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import json
import time
import re
import tempfile
import shutil

# Import our utilities
import sys
sys.path.append(str(Path(__file__).parent.parent))
from utils.configs import OUTPUT_DIR, LOGS_DIR
from utils.error_handler import ErrorHandler
from utils.text_cleaner import TextCleaner

try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

try:
    import win32com.client
    _WIN32COM_AVAILABLE = True
except ImportError:
    _WIN32COM_AVAILABLE = False


class PPTXExtractor:
    """Extract text from PowerPoint files"""
    
    def __init__(self):
        if not _PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )

        self.text_cleaner = TextCleaner()
        self.base_output_dir = OUTPUT_DIR
        self.base_logs_dir = LOGS_DIR

        # Ensure base directories exist
        self.base_output_dir.mkdir(parents=True, exist_ok=True)
        self.base_logs_dir.mkdir(parents=True, exist_ok=True)
    
    def _convert_ppt_to_pptx(self, ppt_path: Path) -> Optional[Path]:
        """
        Convert a legacy .ppt file to .pptx using PowerPoint COM (Windows)
        or LibreOffice as a fallback.

        Returns:
            Path to a temporary .pptx file, or None if conversion failed.
            Caller is responsible for deleting the temp file.
        """
        tmp_dir = Path(tempfile.mkdtemp())
        pptx_path = tmp_dir / (ppt_path.stem + ".pptx")

        # --- Try win32com (requires Microsoft Office installed) ---
        if _WIN32COM_AVAILABLE:
            try:
                import pythoncom
                pythoncom.CoInitialize()
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                ppt_app.Visible = 1
                presentation = ppt_app.Presentations.Open(
                    str(ppt_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False
                )
                # 24 = ppSaveAsOpenXMLPresentation (.pptx)
                presentation.SaveAs(str(pptx_path.resolve()), 24)
                presentation.Close()
                ppt_app.Quit()
                if pptx_path.exists():
                    return pptx_path
            except Exception:
                pass
            finally:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

        # --- Fallback: LibreOffice CLI ---
        import subprocess
        for libreoffice_cmd in ("libreoffice", "soffice"):
            try:
                result = subprocess.run(
                    [libreoffice_cmd, "--headless", "--convert-to", "pptx",
                     "--outdir", str(tmp_dir), str(ppt_path.resolve())],
                    capture_output=True, timeout=60
                )
                if result.returncode == 0 and pptx_path.exists():
                    return pptx_path
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue

        # Cleanup temp dir if nothing worked
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return None

    def _create_resource_name(self, filename: str) -> str:
        """
        Create a clean resource name from filename
        
        Example: "Lecture 5 - AI Basics.pptx" → "lecture_5_ai_basics"
        """
        # Remove extension
        name = Path(filename).stem
        
        # Convert to lowercase
        name = name.lower()
        
        # Replace spaces and special chars with underscore
        name = re.sub(r'[^\w\s-]', '', name)
        name = re.sub(r'[-\s]+', '_', name)
        
        # Remove leading/trailing underscores
        name = name.strip('_')
        
        # Limit length
        if len(name) > 50:
            name = name[:50]
        
        return name or "unnamed_resource"
    
    def _setup_resource_directories(self, resource_name: str, output_dir_override: Optional[Path] = None) -> tuple:
        """
        Create directories for a specific resource

        Args:
            resource_name: Clean name derived from the input filename.
            output_dir_override: If provided, use this directory for output
                instead of creating a resource-specific subfolder.

        Returns:
            (output_dir, logs_dir) paths
        """
        if output_dir_override:
            resource_output_dir = Path(output_dir_override)
        else:
            resource_output_dir = self.base_output_dir / resource_name

        resource_logs_dir = self.base_logs_dir / resource_name

        resource_output_dir.mkdir(parents=True, exist_ok=True)
        resource_logs_dir.mkdir(parents=True, exist_ok=True)

        return resource_output_dir, resource_logs_dir
    
    def extract(self,
                pptx_path: str,
                resource_id: Optional[str] = None,
                clean_text: bool = True,
                include_notes: bool = True,
                output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        Extract text from a PowerPoint file

        Args:
            pptx_path: Path to PPTX file
            resource_id: Optional unique identifier (if None, uses filename)
            clean_text: Whether to clean the extracted text
            include_notes: Whether to include speaker notes
            output_dir: Optional shared output directory. When provided,
                all output files are written here instead of a per-resource subfolder.

        Returns:
            Dictionary with extraction results and metadata

        Example:
            extractor = PPTXExtractor()
            result = extractor.extract("slides.pptx")
            # Creates: output/slides/text.txt and output/slides/metadata.json
        """
        start_time = time.time()
        pptx_path = Path(pptx_path)
        _tmp_dir_to_cleanup = None  # track temp dir for .ppt conversion

        # Auto-convert legacy .ppt to .pptx
        if pptx_path.suffix.lower() == ".ppt":
            converted = self._convert_ppt_to_pptx(pptx_path)
            if converted is None:
                resource_name = self._create_resource_name(pptx_path.name)
                override = Path(output_dir) if output_dir else None
                out_dir, _ = self._setup_resource_directories(resource_name, output_dir_override=override)
                return self._create_error_result(
                    resource_name,
                    "Cannot convert .ppt file: install Microsoft Office (win32com) or LibreOffice.",
                    out_dir,
                    pptx_path.name
                )
            _tmp_dir_to_cleanup = converted.parent
            pptx_path = converted

        # Create resource name from filename
        resource_name = self._create_resource_name(pptx_path.name)

        # Setup directories
        override = Path(output_dir) if output_dir else None
        output_dir, logs_dir = self._setup_resource_directories(resource_name, output_dir_override=override)
        
        # Initialize error handler for this specific resource
        error_handler = ErrorHandler(f"pptx_{resource_name}")
        # Move log file to resource-specific directory
        error_handler.log_file = logs_dir / "extraction.log"
        error_handler.logger = error_handler._setup_logger()
        
        # Validate file
        if not pptx_path.exists():
            error_msg = f"PPTX file not found: {pptx_path}"
            error_handler.log_error(
                FileNotFoundError(error_msg),
                context="Validating PPTX file",
                metadata={"path": str(pptx_path)}
            )
            return self._create_error_result(
                resource_name, error_msg, output_dir, pptx_path.name
            )
        
        # Check file size
        file_size = pptx_path.stat().st_size
        file_size_mb = file_size / (1024 * 1024)
        
        error_handler.log_info(
            f"Starting PPTX extraction: {pptx_path.name}",
            metadata={
                "size_mb": f"{file_size_mb:.2f}",
                "resource_name": resource_name,
                "output_dir": str(output_dir)
            }
        )
        
        try:
            # Open PowerPoint presentation
            prs = Presentation(pptx_path)
            
            # Extract text from all slides
            extracted_text = ""
            slide_count = len(prs.slides)
            slides_with_notes = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide header
                extracted_text += f"\n{'='*60}\n"
                extracted_text += f"SLIDE {slide_num}\n"
                extracted_text += f"{'='*60}\n\n"
                
                # Extract text from shapes (text boxes, titles, content)
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    extracted_text += '\n'.join(slide_text)
                    extracted_text += "\n"
                else:
                    extracted_text += "[No text content on this slide]\n"
                
                # Extract speaker notes if requested
                if include_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        extracted_text += f"\n[SPEAKER NOTES]\n"
                        extracted_text += notes_text + "\n"
                        slides_with_notes += 1
                
                extracted_text += "\n"
            
            # Clean text if requested
            if clean_text:
                extracted_text = self.text_cleaner.clean_text(
                    extracted_text,
                    remove_urls=False,  # Keep URLs
                    remove_emails=False,  # Keep emails
                    fix_spacing=True
                )
            
            # Remove duplicate lines (slides often have duplicates)
            extracted_text = self.text_cleaner.remove_duplicate_lines(extracted_text)
            
            # Calculate processing time
            processing_time = time.time() - start_time
            
            # Create metadata
            metadata = {
                "resource_name": resource_name,
                "resource_id": resource_id or resource_name,
                "filename": pptx_path.name,
                "source_type": "pptx",
                "upload_date": datetime.now().isoformat(),
                "extraction_timestamp": datetime.now().isoformat(),
                "file_size_bytes": file_size,
                "processing_time_seconds": round(processing_time, 2),
                "status": "success",
                "error_message": None,
                "slide_count": slide_count,
                "slides_with_notes": slides_with_notes,
                "character_count": len(extracted_text),
                "included_notes": include_notes
            }
            
            # Save text file
            text_file = output_dir / "text.txt"
            text_file.write_text(extracted_text, encoding='utf-8')
            metadata["extracted_text_path"] = str(text_file)
            
            # Save metadata file
            metadata_file = output_dir / "metadata.json"
            metadata_file.write_text(json.dumps(metadata, indent=2), encoding='utf-8')
            
            error_handler.log_success(
                f"PPTX extracted successfully: {pptx_path.name}",
                metadata={
                    "slides": slide_count,
                    "chars": len(extracted_text),
                    "time": f"{processing_time:.2f}s",
                    "output": str(output_dir)
                }
            )
            
            return {
                "success": True,
                "resource_name": resource_name,
                "resource_id": resource_id or resource_name,
                "text_file": str(text_file),
                "metadata_file": str(metadata_file),
                "output_dir": str(output_dir),
                "logs_dir": str(logs_dir),
                "metadata": metadata,
                "extracted_text": extracted_text
            }
            
        except Exception as e:
            processing_time = time.time() - start_time

            error_handler.log_error(
                e,
                context=f"Extracting PPTX: {pptx_path.name}",
                metadata={"resource_name": resource_name}
            )

            return self._create_error_result(
                resource_name,
                str(e),
                output_dir,
                pptx_path.name,
                file_size,
                processing_time
            )
        finally:
            if _tmp_dir_to_cleanup:
                shutil.rmtree(_tmp_dir_to_cleanup, ignore_errors=True)
    
    def _create_error_result(self,
                           resource_name: str,
                           error_message: str,
                           output_dir: Path,
                           filename: str = "unknown",
                           file_size: int = 0,
                           processing_time: float = 0) -> Dict[str, Any]:
        """Create error result when extraction fails"""
        
        metadata = {
            "resource_name": resource_name,
            "filename": filename,
            "source_type": "pptx",
            "upload_date": datetime.now().isoformat(),
            "extraction_timestamp": datetime.now().isoformat(),
            "file_size_bytes": file_size,
            "processing_time_seconds": round(processing_time, 2),
            "status": "failed",
            "error_message": error_message
        }
        
        # Save metadata even for failed extraction
        metadata_file = output_dir / "metadata.json"
        metadata_file.write_text(json.dumps(metadata, indent=2), encoding='utf-8')
        
        return {
            "success": False,
            "resource_name": resource_name,
            "text_file": None,
            "metadata_file": str(metadata_file),
            "output_dir": str(output_dir),
            "metadata": metadata,
            "extracted_text": "",
            "error": error_message
        }
    
    def extract_metadata_only(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract only metadata without extracting text
        Useful for quick file info
        
        Returns:
            Dictionary with PPTX metadata
        """
        _tmp_dir = None
        try:
            pptx_path = Path(pptx_path)
            if pptx_path.suffix.lower() == ".ppt":
                converted = self._convert_ppt_to_pptx(pptx_path)
                if converted is None:
                    return {}
                _tmp_dir = converted.parent
                pptx_path = converted
            prs = Presentation(pptx_path)
            
            # Count slides with notes
            slides_with_notes = sum(1 for slide in prs.slides if slide.has_notes_slide)
            
            metadata = {
                "filename": pptx_path.name,
                "resource_name": self._create_resource_name(pptx_path.name),
                "slide_count": len(prs.slides),
                "slides_with_notes": slides_with_notes,
                "file_size_bytes": pptx_path.stat().st_size,
                "file_size_mb": round(pptx_path.stat().st_size / (1024 * 1024), 2),
                "core_properties": {
                    "title": prs.core_properties.title or "N/A",
                    "author": prs.core_properties.author or "N/A",
                    "subject": prs.core_properties.subject or "N/A",
                    "created": str(prs.core_properties.created) if prs.core_properties.created else "N/A",
                    "modified": str(prs.core_properties.modified) if prs.core_properties.modified else "N/A"
                }
            }
            
            return metadata

        except Exception as e:
            error_handler = ErrorHandler("pptx_metadata")
            error_handler.log_error(
                e,
                context=f"Extracting metadata from {pptx_path}"
            )
            return {}
        finally:
            if _tmp_dir:
                shutil.rmtree(_tmp_dir, ignore_errors=True)


# Example usage and testing
if __name__ == "__main__":
    from utils.file_picker import FilePicker
    
    print("=== Testing PPTX Extractor ===\n")
    
    # Initialize extractor
    extractor = PPTXExtractor()
    
    # Use file picker to select PPTX
    picker = FilePicker()
    print("Please select a PowerPoint file...")
    test_pptx = picker.pick_pptx()
    picker.close()
    
    if test_pptx:
        print(f"\n✓ Selected: {Path(test_pptx).name}\n")
        
        print("1. Extracting metadata only...")
        metadata = extractor.extract_metadata_only(test_pptx)
        print(f"   Resource name: {metadata.get('resource_name', 'N/A')}")
        print(f"   Slides: {metadata.get('slide_count', 'N/A')}")
        print(f"   Slides with notes: {metadata.get('slides_with_notes', 'N/A')}")
        print(f"   Size: {metadata.get('file_size_mb', 'N/A')} MB")
        print(f"   Title: {metadata.get('core_properties', {}).get('title', 'N/A')}")
        print(f"   Author: {metadata.get('core_properties', {}).get('author', 'N/A')}\n")
        
        print("2. Full extraction (with speaker notes)...")
        result = extractor.extract(
            pptx_path=test_pptx,
            clean_text=True,
            include_notes=True
        )
        
        if result['success']:
            print(f"   ✓ Success!")
            print(f"   Resource name: {result['resource_name']}")
            print(f"   Output directory: {result['output_dir']}")
            print(f"   Logs directory: {result['logs_dir']}")
            print(f"   Text file: {result['text_file']}")
            print(f"   Metadata file: {result['metadata_file']}")
            print(f"   Slides: {result['metadata']['slide_count']}")
            print(f"   Slides with notes: {result['metadata']['slides_with_notes']}")
            print(f"   Characters: {result['metadata']['character_count']}")
            print(f"   Processing time: {result['metadata']['processing_time_seconds']}s")
            
            print(f"\n   Preview (first 400 chars):")
            print("   " + "-" * 50)
            preview = result['extracted_text'][:400]
            print(f"   {preview}...")
            print("   " + "-" * 50)
        else:
            print(f"   ✗ Failed: {result['error']}")
    else:
        print("❌ No file selected")
        print("   The extractor is ready to use!")
